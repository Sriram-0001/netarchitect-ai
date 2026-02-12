import streamlit as st
import plotly.graph_objects as go
from agents.analysis_engine import run_analysis
from pptx import Presentation
from gtts import gTTS
import tempfile
import os
import base64

# ---------------------------------------------------------
# PAGE CONFIG
# ---------------------------------------------------------
st.set_page_config(page_title="NetArchitect AI", layout="wide")

# ---------------------------------------------------------
# PREMIUM EXECUTIVE THEME
# ---------------------------------------------------------
st.markdown("""
<style>
body { background: linear-gradient(135deg,#0f172a,#020617); color:white; }

.kpi-card {
    background: linear-gradient(145deg,#1e293b,#0f172a);
    padding:30px;
    border-radius:20px;
    text-align:center;
    box-shadow: 0 20px 60px rgba(0,255,150,0.15);
    transition: all 0.3s ease;
}
.kpi-card:hover {
    transform: translateY(-8px) scale(1.05);
    box-shadow: 0 25px 80px rgba(16,185,129,0.4);
}

.kpi-title { color:#94a3b8; font-size:18px; }
.kpi-value { font-size:36px; font-weight:800; color:#10b981; }

.section-title {
    font-size:28px;
    font-weight:700;
    margin-top:40px;
    margin-bottom:20px;
    color:#22d3ee;
}

.verdict-box {
    background: linear-gradient(145deg,#111827,#1e293b);
    padding:35px;
    border-radius:20px;
    border-left: 6px solid #22c55e;
    box-shadow: 0 20px 60px rgba(0,255,150,0.2);
}

.ai-avatar {
    font-size:60px;
    text-align:center;
    animation: pulse 1.5s infinite;
}

@keyframes pulse {
  0% { transform: scale(1); opacity:0.8; }
  50% { transform: scale(1.1); opacity:1; }
  100% { transform: scale(1); opacity:0.8; }
}
</style>
""", unsafe_allow_html=True)

# ---------------------------------------------------------
# HEADER
# ---------------------------------------------------------
st.title("🚀 NetArchitect AI")
st.caption("Enterprise Infrastructure Intelligence Platform")

# ---------------------------------------------------------
# INPUT CONTROLS
# ---------------------------------------------------------
col1, col2, col3 = st.columns(3)

with col1:
    employees = st.slider("Employees", 10, 500, 120)
    routers = st.slider("Routers", 1, 5, 1)

with col2:
    switches = st.slider("Switches", 1, 10, 3)
    aps = st.slider("Access Points", 1, 10, 4)

with col3:
    firewalls = st.slider("Firewalls", 1, 3, 1)
    ids = st.slider("IDS Systems", 0, 3, 1)

budget = st.slider("Budget", 100000, 2000000, 700000, step=50000)

cloud_model = st.selectbox("Deployment Model", ["Cloud", "Hybrid", "On-Prem"])

# ---------------------------------------------------------
# BUILD INFRA PACKAGE
# ---------------------------------------------------------
infra_package = {
    "company_profile": {
        "num_employees": employees,
        "office_size_sqft": 6000,
        "security_level": "High",
        "growth_rate_percent": 20,
        "cloud_preference": cloud_model,
        "budget": budget
    },
    "infrastructure_design": {
        "topology": "Hybrid",
        "components": {
            "routers": routers,
            "switches": switches,
            "access_points": aps,
            "firewalls": firewalls,
            "ids_systems": ids
        },
        "cloud_architecture": {
            "model": cloud_model,
            "cloud_servers": 2,
            "on_prem_servers": 2
        },
        "redundancy": {"enabled": True, "dual_isp": True},
        "selected_models": {
            "cisco": {
                "router_model": "Cisco Catalyst 8300",
                "switch_model": "Cisco Catalyst 9200",
                "access_point_model": "Cisco Catalyst 9120",
                "firewall_model": "Cisco Firepower 1120"
            },
            "tplink": {
                "router_model": "TP-Link ER8411",
                "switch_model": "TP-Link JetStream 48",
                "access_point_model": "TP-Link EAP660 HD",
                "firewall_model": "TP-Link SafeStream TL-R605"
            }
        }
    },
    "scalability_projection": {
        "year_1_users": 144,
        "year_2_users": 172,
        "year_3_users": 207,
        "upgrade_required": True,
        "upgrade_reason": "Projected growth"
    }
}

# ---------------------------------------------------------
# OPTIMIZATION BUTTON
# ---------------------------------------------------------
if st.button("🚀 Run Optimization"):

    result = run_analysis(infra_package)

    cost = result["cost_analysis"]
    performance = result["performance_scores"]
    risk = result["security_risk_scores"]
    recommendation = result["optimization_recommendation"]

    # ---------------------------------------------------------
    # KPI CARDS
    # ---------------------------------------------------------
    colA, colB, colC = st.columns(3)

    colA.markdown(f"""
    <div class="kpi-card">
        <div class="kpi-title">Recommended Vendor</div>
        <div class="kpi-value">{recommendation['recommended_vendor'].upper()}</div>
    </div>
    """, unsafe_allow_html=True)

    colB.markdown(f"""
    <div class="kpi-card">
        <div class="kpi-title">Cost Difference</div>
        <div class="kpi-value">₹ {cost['cost_difference']:,.0f}</div>
    </div>
    """, unsafe_allow_html=True)

    colC.markdown(f"""
    <div class="kpi-card">
        <div class="kpi-title">Optimization Required</div>
        <div class="kpi-value">{recommendation['optimization_needed']}</div>
    </div>
    """, unsafe_allow_html=True)

    # ---------------------------------------------------------
    # PERFORMANCE GAUGE
    # ---------------------------------------------------------
    st.markdown('<div class="section-title">Performance Gauge</div>', unsafe_allow_html=True)

    gauge_fig = go.Figure(go.Indicator(
        mode="gauge+number",
        value=performance["cisco"],
        gauge={
            'axis': {'range': [0,100]},
            'bar': {'color': "#10b981"},
            'steps': [
                {'range': [0,40], 'color': "#ef4444"},
                {'range': [40,70], 'color': "#f59e0b"},
                {'range': [70,100], 'color': "#10b981"}
            ]
        }
    ))

    gauge_fig.update_layout(height=300)
    st.plotly_chart(gauge_fig, use_container_width=True)

    # ---------------------------------------------------------
    # FINAL VERDICT
    # ---------------------------------------------------------
    st.markdown('<div class="section-title">Final Decision Verdict</div>', unsafe_allow_html=True)

    verdict_text = f"""
    Recommended Vendor: {recommendation['recommended_vendor'].upper()}
    Reason: {recommendation['reason']}
    Budget Status: {'Within Budget' if not recommendation['optimization_needed'] else 'Exceeds Budget - Optimization Suggested'}

    Cisco Performance: {performance['cisco']}
    TP-Link Performance: {performance['tplink']}

    Cisco Risk: {risk['cisco']}
    TP-Link Risk: {risk['tplink']}
    """

    st.markdown(f"""
    <div class="verdict-box">
    {verdict_text}
    </div>
    """, unsafe_allow_html=True)

    # ---------------------------------------------------------
    # AI EXECUTIVE NARRATION
    # ---------------------------------------------------------
    st.markdown('<div class="section-title">🤖 AI Executive Briefing</div>', unsafe_allow_html=True)

    st.markdown('<div class="ai-avatar">🤖</div>', unsafe_allow_html=True)

    tts = gTTS(verdict_text)
    audio_path = os.path.join(tempfile.gettempdir(), "ai_voice.mp3")
    tts.save(audio_path)

    with open(audio_path, "rb") as f:
        audio_bytes = f.read()
        b64 = base64.b64encode(audio_bytes).decode()

    audio_html = f"""
    <audio autoplay>
    <source src="data:audio/mp3;base64,{b64}" type="audio/mp3">
    </audio>
    """

    st.markdown(audio_html, unsafe_allow_html=True)

    # ---------------------------------------------------------
    # PPT EXPORT
    # ---------------------------------------------------------
    if st.button("📊 Generate Executive Slide Deck"):

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "NetArchitect AI"

        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Final Recommendation"
        slide2.placeholders[1].text = verdict_text

        temp_path = os.path.join(tempfile.gettempdir(), "PitchDeck.pptx")
        prs.save(temp_path)

        with open(temp_path, "rb") as f:
            st.download_button("Download PPT", f,
                               file_name="NetArchitect_AI_Pitch.pptx")
